#!/usr/bin/env python3
"""Generate a single-slide PowerPoint containing a structured challenge table.

Delegates table rendering to hx-pptx's shared table_renderer module.

Usage:
    python generate_table_slide.py data.json -o output.pptx [-t template.pptx]

Requires table_renderer on PYTHONPATH (set by the SKILL.md bash command).
"""

import argparse
import json
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

from table_renderer import render_table

# Colours used for the title/subtitle/footer (table colours live in table_renderer)
DARK_TEXT = RGBColor(0x1C, 0x27, 0x33)
GREY = RGBColor(0x83, 0x8D, 0x99)

SLIDE_WIDTH_CM = 33.87
SLIDE_HEIGHT_CM = 19.05

TABLE_LEFT_CM = 1.0
TABLE_TOP_CM = 3.2
TABLE_WIDTH_CM = 31.87
TABLE_BOTTOM_MARGIN_CM = 1.8

COLUMNS = [
    {"header": "Challenge", "width": 0.11},
    {"header": "Issue", "width": 0.19},
    {"header": "Impact", "width": 0.20},
    {"header": "Measure", "width": 0.18},
    {"header": "Org Context", "width": 0.18},
    {"header": "Objective", "width": 0.14},
]

FIELDS = ["challenge", "issue", "impact", "measure", "org_context", "objective"]

TABLE_POSITION = {
    "left_cm": TABLE_LEFT_CM,
    "top_cm": TABLE_TOP_CM,
    "width_cm": TABLE_WIDTH_CM,
    "height_cm": SLIDE_HEIGHT_CM - TABLE_TOP_CM - TABLE_BOTTOM_MARGIN_CM,
}


def generate_table_slide(data, output_path, template_path=None):
    """Generate a single-slide .pptx with a challenge table."""
    account_name = data["account_name"]
    rows = data["rows"]

    if len(rows) > 5:
        print(
            f"Warning: {len(rows)} rows provided, table may not fit on one slide. "
            "Consider limiting to 5.",
            file=sys.stderr,
        )

    # ── Presentation setup ────────────────────────────────────────────────
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
        sldIdLst = prs.slides._sldIdLst
        for item in list(sldIdLst):
            ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
            rId = item.get(f"{{{ns}}}id")
            if rId:
                try:
                    prs.part.drop_rel(rId)
                except KeyError:
                    pass
            sldIdLst.remove(item)
        layout = None
        for master in prs.slide_masters:
            for sl in master.slide_layouts:
                if "blank" in sl.name.lower() and "white" in sl.name.lower():
                    layout = sl
                    break
            if layout:
                break
        if layout is None:
            layout = prs.slide_layouts[0]
    else:
        prs = Presentation()
        prs.slide_width = Cm(SLIDE_WIDTH_CM)
        prs.slide_height = Cm(SLIDE_HEIGHT_CM)
        layout = prs.slide_layouts[6]

    slide = prs.slides.add_slide(layout)

    # ── Title ─────────────────────────────────────────────────────────────
    txbox = slide.shapes.add_textbox(
        Cm(TABLE_LEFT_CM), Cm(0.8), Cm(TABLE_WIDTH_CM), Cm(1.2),
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{account_name} - Structured Challenge Table"
    p.font.size = Pt(16)
    p.font.name = "Arial"
    p.font.color.rgb = DARK_TEXT
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    subtitle = tf.add_paragraph()
    subtitle.text = (
        "Transformation of strategy-slide content into executive-ready "
        "business challenges. Source: internal account intelligence, "
        "public filings, and executive commentary."
    )
    subtitle.font.size = Pt(7)
    subtitle.font.name = "Arial"
    subtitle.font.color.rgb = GREY
    subtitle.font.italic = True
    subtitle.alignment = PP_ALIGN.LEFT

    # ── Challenge table (delegated to table_renderer) ─────────────────────
    table_def = {
        "columns": COLUMNS,
        "rows": [[r.get(f, "") for f in FIELDS] for r in rows],
        "style": {"bold_first_col": True},
    }
    render_table(slide, table_def, colour="Green", position=TABLE_POSITION)

    # ── Footer ────────────────────────────────────────────────────────────
    footer_box = slide.shapes.add_textbox(
        Cm(TABLE_LEFT_CM), Cm(SLIDE_HEIGHT_CM - 1.2),
        Cm(TABLE_WIDTH_CM), Cm(0.8),
    )
    footer_p = footer_box.text_frame.paragraphs[0]
    footer_p.text = f"Prepared by hyperexponential - {datetime.now().strftime('%B %Y')}"
    footer_p.font.size = Pt(7)
    footer_p.font.name = "Arial"
    footer_p.font.color.rgb = GREY
    footer_p.font.italic = True
    footer_p.alignment = PP_ALIGN.RIGHT

    prs.save(output_path)
    print(f"Generated challenge table ({len(rows)} rows) -> {output_path}")


def main():
    parser = argparse.ArgumentParser(
        description="Generate a challenge table PowerPoint slide",
    )
    parser.add_argument("data", help="Path to JSON data file")
    parser.add_argument("-o", "--output", required=True, help="Output .pptx path")
    parser.add_argument("-t", "--template", help="Optional hx template .pptx")
    args = parser.parse_args()

    if not Path(args.data).exists():
        print(f"Error: Data file not found: {args.data}", file=sys.stderr)
        sys.exit(1)

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    if "account_name" not in data:
        print("Error: JSON must include 'account_name'", file=sys.stderr)
        sys.exit(1)
    if "rows" not in data or not data["rows"]:
        print("Error: JSON must include 'rows' with at least one entry", file=sys.stderr)
        sys.exit(1)

    generate_table_slide(data, args.output, template_path=args.template)


if __name__ == "__main__":
    main()
